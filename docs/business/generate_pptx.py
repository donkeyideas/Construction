"""
Generate Buildwrk Pitch Deck PowerPoint
Matches the homepage design system:
  - Primary Blue: #1d4ed8
  - Amber Accent: #b45309
  - Background: #ffffff / #f5f0eb
  - Text: #292524
  - Fonts: Playfair Display (headings), Inter (body) -> mapped to Calibri/Georgia
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Brand Colors ───
BLUE = RGBColor(0x1D, 0x4E, 0xD8)       # #1d4ed8
BLUE_DARK = RGBColor(0x1E, 0x40, 0xAF)   # #1e40af
AMBER = RGBColor(0xB4, 0x53, 0x09)       # #b45309
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_WARM = RGBColor(0xF5, 0xF0, 0xEB)     # #f5f0eb
TEXT_DARK = RGBColor(0x29, 0x25, 0x24)    # #292524
MUTED = RGBColor(0x78, 0x71, 0x6C)       # #78716c
GREEN = RGBColor(0x16, 0xA3, 0x4A)       # #16a34a
RED = RGBColor(0xDC, 0x26, 0x26)         # #dc2626
LIGHT_BLUE_BG = RGBColor(0xEB, 0xF1, 0xFD) # light blue tint

# ─── Fonts ───
HEADING_FONT = "Georgia"    # Closest system font to Playfair Display
BODY_FONT = "Calibri"       # Closest system font to Inter

# ─── Presentation Setup ───
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def add_bg(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_name=BODY_FONT, color=TEXT_DARK, bold=False,
                 alignment=PP_ALIGN.LEFT, line_spacing=1.2):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(0)
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=16,
                   default_color=TEXT_DARK, default_font=BODY_FONT, line_spacing=1.5):
    """Add a text box with multiple styled paragraphs.
    lines = [(text, font_size, color, bold, font_name), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = line[0] if len(line) > 0 else ""
        size = line[1] if len(line) > 1 else default_size
        color = line[2] if len(line) > 2 else default_color
        bold = line[3] if len(line) > 3 else False
        font = line[4] if len(line) > 4 else default_font
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.space_after = Pt(int(size * 0.4))
    return txBox

def add_amber_bar(slide, left, top, width=Inches(0.8), height=Pt(3)):
    """Add the signature amber accent bar (matches homepage hero)."""
    return add_shape(slide, left, top, width, height, AMBER)

def add_blue_header_bar(slide):
    """Full-width blue header bar at top of slide."""
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE)

def add_slide_number(slide, num, total=14):
    """Add slide number in bottom right."""
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)

def add_logo(slide, left=Inches(0.6), top=Inches(0.2)):
    """Add Buildwrk text logo."""
    add_text_box(slide, left, top, Inches(2.5), Inches(0.5),
                 "Buildwrk", font_size=20, font_name=HEADING_FONT,
                 color=BLUE, bold=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG_WARM)

# Blue accent block on left
add_shape(slide, Inches(0), Inches(0), Inches(5.5), SLIDE_H, BLUE)

# Title text on blue area
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(4.2), Inches(1.0),
             "Buildwrk", font_size=48, font_name=HEADING_FONT, color=WHITE, bold=True)

add_amber_bar(slide, Inches(0.8), Inches(2.6), Inches(1.2), Pt(4))

add_multi_text(slide, Inches(0.8), Inches(2.8), Inches(4.2), Inches(2.5), [
    ("The Unified Construction", 28, WHITE, False, HEADING_FONT),
    ("Intelligence Platform", 28, WHITE, False, HEADING_FONT),
    ("", 12, WHITE),
    ("One platform to manage projects, finances,", 14, RGBColor(0xD1, 0xD5, 0xDB), False),
    ("safety, and operations for the $1.3T", 14, RGBColor(0xD1, 0xD5, 0xDB), False),
    ("construction industry.", 14, RGBColor(0xD1, 0xD5, 0xDB), False),
])

# Right side - key details on warm bg
add_multi_text(slide, Inches(6.2), Inches(2.0), Inches(6.5), Inches(4.5), [
    ("Pre-Seed Round", 32, TEXT_DARK, True, HEADING_FONT),
    ("", 10, TEXT_DARK),
    ("$750K Target Raise", 24, BLUE, True),
    ("$4M-$6M Pre-Money Valuation", 18, MUTED),
    ("", 10, TEXT_DARK),
    ("February 2026", 16, AMBER, True),
    ("", 20, TEXT_DARK),
    ("80,000+ lines of production code", 14, TEXT_DARK),
    ("19 integrated modules  |  46 database tables", 14, TEXT_DARK),
    ("461 security policies  |  9 AI providers", 14, TEXT_DARK),
    ("13 demo datasets  |  Live in production", 14, TEXT_DARK),
])

add_slide_number(slide, 1)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_header_bar(slide)
add_logo(slide)

add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
             "Construction companies run on duct tape",
             font_size=32, font_name=HEADING_FONT, color=TEXT_DARK, bold=True)

add_amber_bar(slide, Inches(0.6), Inches(1.5), Inches(1.0))

add_text_box(slide, Inches(0.6), Inches(1.8), Inches(11), Inches(0.6),
             "A typical mid-market GC uses 5-8 disconnected software tools costing $150K-$378K/year:",
             font_size=16, color=MUTED)

# Software stack table - boxes
tools = [
    ("Procore", "$60K-$112K/yr", "Project management only"),
    ("Sage 300 CRE", "$36K-$96K/yr", "Accounting only"),
    ("SafetyCulture", "$6K-$24K/yr", "Safety only"),
    ("Box/SharePoint", "$6K-$18K/yr", "Documents only"),
    ("ADP/Paychex", "$12K-$36K/yr", "Payroll only"),
]

for i, (name, cost, desc) in enumerate(tools):
    left = Inches(0.6 + (i * 2.45))
    top = Inches(2.5)

    # Card background
    card = add_shape(slide, left, top, Inches(2.2), Inches(1.8), BG_WARM)
    card.line.color.rgb = RGBColor(0xD6, 0xD3, 0xD1)
    card.line.width = Pt(1)

    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), Inches(1.9), Inches(0.35),
                 name, font_size=16, color=TEXT_DARK, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.55), Inches(1.9), Inches(0.35),
                 cost, font_size=20, color=RED, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(1.1), Inches(1.9), Inches(0.5),
                 desc, font_size=12, color=MUTED)

# Bottom impact stats
add_text_box(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5),
             "The Result:", font_size=20, color=TEXT_DARK, bold=True)

impacts = [
    "Change orders in Procore require manual journal entries in Sage",
    "Safety incidents don't link to project budgets for insurance accruals",
    "5-10% of project revenue lost to data reconciliation overhead",
]
for i, text in enumerate(impacts):
    add_text_box(slide, Inches(0.8), Inches(5.4 + i * 0.45), Inches(11), Inches(0.4),
                 f"\u2022  {text}", font_size=14, color=MUTED)

add_slide_number(slide, 2)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: The Solution
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_header_bar(slide)
add_logo(slide)

add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
             "Buildwrk: Everything in one platform",
             font_size=32, font_name=HEADING_FONT, color=TEXT_DARK, bold=True)

add_amber_bar(slide, Inches(0.6), Inches(1.5), Inches(1.0))

# Module grid - 2 columns
modules_left = [
    ("Project Management", "Gantt, RFIs, submittals, daily logs, change orders", "Replaces Procore"),
    ("Financial Accounting", "GL, AR/AP, job costing, financial statements", "Replaces Sage/QuickBooks"),
    ("Safety Compliance", "Incidents, inspections, toolbox talks, OSHA", "Replaces SafetyCulture"),
    ("Equipment Management", "Inventory, assignments, depreciation, maintenance", "Replaces HCSS"),
]

modules_right = [
    ("Property Management", "Leases, units, maintenance, rent, NOI tracking", "Replaces Yardi"),
    ("CRM & Bidding", "Pipeline, proposals, bid tracking, analytics", "Replaces Salesforce"),
    ("AI Analytics", "9 LLM providers, natural-language data queries", "Nothing exists like it"),
    ("Mobile PWA", "Clock in/out, daily logs, photos, tasks from field", "Replaces ExakTime"),
]

for col, modules in enumerate([modules_left, modules_right]):
    for i, (name, desc, replaces) in enumerate(modules):
        left = Inches(0.6 + col * 6.3)
        top = Inches(1.9 + i * 1.25)

        # Blue number circle
        num_shape = add_shape(slide, left, top + Inches(0.05), Inches(0.35), Inches(0.35), BLUE)
        add_text_box(slide, left + Inches(0.02), top + Inches(0.05), Inches(0.35), Inches(0.35),
                     str(i + 1 + col * 4), font_size=14, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, left + Inches(0.5), top, Inches(5.5), Inches(0.3),
                     name, font_size=16, color=TEXT_DARK, bold=True)
        add_text_box(slide, left + Inches(0.5), top + Inches(0.3), Inches(5.5), Inches(0.3),
                     desc, font_size=12, color=MUTED)
        add_text_box(slide, left + Inches(0.5), top + Inches(0.6), Inches(5.5), Inches(0.3),
                     replaces, font_size=11, color=AMBER, bold=True)

add_slide_number(slide, 3)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: How It Works (Automation)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_header_bar(slide)
add_logo(slide)

add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
             "When a change order is approved...",
             font_size=32, font_name=HEADING_FONT, color=TEXT_DARK, bold=True)

add_amber_bar(slide, Inches(0.6), Inches(1.5), Inches(1.0))

# Flow steps
steps = [
    ("Budget updated", "Project budget reflects new contract value"),
    ("Journal entry created", "DR Accounts Receivable / CR Revenue (balanced)"),
    ("Contract adjusted", "Total contract value updated automatically"),
    ("Audit trail logged", "Who approved, when, original vs. new amount"),
    ("Team notified", "PM and accountant receive instant alerts"),
]

for i, (title, desc) in enumerate(steps):
    top = Inches(2.0 + i * 0.95)

    # Step number in blue circle
    add_shape(slide, Inches(1.5), top, Inches(0.5), Inches(0.5), BLUE)
    add_text_box(slide, Inches(1.5), top + Inches(0.02), Inches(0.5), Inches(0.5),
                 str(i + 1), font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Connecting line to next step
    if i < len(steps) - 1:
        add_shape(slide, Inches(1.72), top + Inches(0.5), Pt(2), Inches(0.45), RGBColor(0xD6, 0xD3, 0xD1))

    # Step text
    add_text_box(slide, Inches(2.3), top + Inches(0.0), Inches(5), Inches(0.3),
                 title, font_size=18, color=TEXT_DARK, bold=True)
    add_text_box(slide, Inches(2.3), top + Inches(0.3), Inches(5), Inches(0.3),
                 desc, font_size=13, color=MUTED)

# Right side - stats box
stat_box = add_shape(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(4.5), BG_WARM)
stat_box.line.fill.background()

add_text_box(slide, Inches(8.8), Inches(2.2), Inches(3.8), Inches(0.4),
             "Zero manual entry", font_size=22, color=BLUE, bold=True, font_name=HEADING_FONT)
add_text_box(slide, Inches(8.8), Inches(2.7), Inches(3.8), Inches(0.4),
             "Zero reconciliation", font_size=22, color=BLUE, bold=True, font_name=HEADING_FONT)

stats = [
    ("5 seconds", "vs 30-60 minutes manually"),
    ("0 errors", "vs 6 data entry points"),
    ("1 click", "vs 6 separate systems"),
    ("Real-time", "vs end-of-month discovery"),
]

for i, (stat, desc) in enumerate(stats):
    top = Inches(3.4 + i * 0.7)
    add_text_box(slide, Inches(8.8), top, Inches(3.5), Inches(0.3),
                 stat, font_size=18, color=TEXT_DARK, bold=True)
    add_text_box(slide, Inches(8.8), top + Inches(0.3), Inches(3.5), Inches(0.3),
                 desc, font_size=12, color=MUTED)

add_slide_number(slide, 4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5: Market Opportunity
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_header_bar(slide)
add_logo(slide)

add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
             "$11.6B market growing at 8.9% CAGR",
             font_size=32, font_name=HEADING_FONT, color=TEXT_DARK, bold=True)

add_amber_bar(slide, Inches(0.6), Inches(1.5), Inches(1.0))

# TAM/SAM/SOM circles
for i, (label, amount, size_in) in enumerate([
    ("TAM", "$11.6B", 3.5), ("SAM", "$1.24B", 2.8), ("SOM (Yr 3)", "$3.3M", 2.0)
]):
    cx = Inches(2.5 + i * 3.5)
    cy = Inches(3.2)
    dim = Inches(size_in)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - dim/2, cy - dim/2, dim, dim
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = [BLUE, RGBColor(0x3B, 0x82, 0xF6), AMBER][i]
    circle.line.fill.background()

    add_text_box(slide, cx - Inches(1), cy - Inches(0.5), Inches(2), Inches(0.4),
                 label, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, cx - Inches(1), cy - Inches(0.1), Inches(2), Inches(0.5),
                 amount, font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 font_name=HEADING_FONT)

# Market segments below
add_text_box(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
             "Our Beachhead: The Underserved Mid-Market", font_size=20, color=TEXT_DARK, bold=True)

segments = [
    ("Small GCs ($1M-$10M)", "650,000+ firms", "$3K-$12K/yr current", "$1.2K-$3K Buildwrk"),
    ("Mid-Market ($10M-$200M)", "45,000+ firms", "$50K-$200K/yr current", "$3K-$6K Buildwrk"),
    ("Large ($200M+)", "5,000+ firms", "$200K-$500K/yr current", "$6K+ Buildwrk"),
]

for i, (seg, firms, current, ours) in enumerate(segments):
    left = Inches(0.6 + i * 4.2)
    top = Inches(5.7)
    card = add_shape(slide, left, top, Inches(3.8), Inches(1.4), BG_WARM)
    card.line.fill.background()
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(3.5), Inches(0.3),
                 seg, font_size=14, color=TEXT_DARK, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.4), Inches(3.5), Inches(0.25),
                 firms, font_size=12, color=MUTED)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.65), Inches(3.5), Inches(0.25),
                 current, font_size=12, color=RED)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.9), Inches(3.5), Inches(0.3),
                 ours, font_size=14, color=GREEN, bold=True)

add_slide_number(slide, 5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: Competitive Advantage
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_header_bar(slide)
add_logo(slide)

add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
             "What we have that others don't",
             font_size=32, font_name=HEADING_FONT, color=TEXT_DARK, bold=True)

add_amber_bar(slide, Inches(0.6), Inches(1.5), Inches(1.0))

# Comparison table header
cols = [("Capability", 3.0), ("Procore", 1.8), ("Sage", 1.5), ("Buildertrend", 2.0), ("Buildwrk", 2.0)]
header_left = Inches(0.6)
header_top = Inches(1.9)

x = header_left
for col_name, col_w in cols:
    bg_color = BLUE if col_name == "Buildwrk" else TEXT_DARK
    add_shape(slide, x, header_top, Inches(col_w), Inches(0.45), bg_color)
    add_text_box(slide, x + Inches(0.1), header_top + Inches(0.02), Inches(col_w - 0.2), Inches(0.4),
                 col_name, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    x += Inches(col_w) + Inches(0.1)

# Table rows
rows = [
    ("Project Management", "Yes", "No", "Yes", "Yes"),
    ("GAAP Accounting", "No", "Yes", "Basic", "Yes"),
    ("Property Management", "No", "No", "No", "Yes"),
    ("Safety Compliance", "Basic", "No", "No", "Yes"),
    ("AI Analytics (9 providers)", "No", "No", "No", "Yes"),
    ("Multi-tenant Portals", "Basic", "No", "No", "Yes"),
    ("Mobile PWA", "Yes", "No", "Yes", "Yes"),
    ("Annual Cost (25 users)", "$60K-$112K", "$36K-$96K", "$6K", "$3K-$6K"),
]

for i, (cap, p, s, b, bw) in enumerate(rows):
    y = header_top + Inches(0.55 + i * 0.45)
    bg = BG_WARM if i % 2 == 0 else WHITE
    x = header_left
    for j, (val, (_, col_w)) in enumerate(zip([cap, p, s, b, bw], cols)):
        add_shape(slide, x, y, Inches(col_w), Inches(0.4), bg)
        if j == 4:  # Buildwrk column
            color = GREEN if val == "Yes" else BLUE
            add_text_box(slide, x + Inches(0.1), y + Inches(0.01), Inches(col_w - 0.2), Inches(0.35),
                         val, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        elif val == "No":
            add_text_box(slide, x + Inches(0.1), y + Inches(0.01), Inches(col_w - 0.2), Inches(0.35),
                         val, font_size=12, color=RED, alignment=PP_ALIGN.CENTER)
        elif val == "Basic":
            add_text_box(slide, x + Inches(0.1), y + Inches(0.01), Inches(col_w - 0.2), Inches(0.35),
                         val, font_size=12, color=AMBER, alignment=PP_ALIGN.CENTER)
        else:
            add_text_box(slide, x + Inches(0.1), y + Inches(0.01), Inches(col_w - 0.2), Inches(0.35),
                         val, font_size=12 if j > 0 else 12, color=TEXT_DARK,
                         bold=(j == 0), alignment=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
        x += Inches(col_w) + Inches(0.1)

add_slide_number(slide, 6)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: Business Model / Pricing
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_header_bar(slide)
add_logo(slide)

add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
             "SaaS subscription with expansion revenue",
             font_size=32, font_name=HEADING_FONT, color=TEXT_DARK, bold=True)

add_amber_bar(slide, Inches(0.6), Inches(1.5), Inches(1.0))

# Three pricing cards
tiers = [
    ("Starter", "$99", "/month", "Up to 5 users", "3 active projects",
     "PM + Safety + Docs", "CSV import", "Email support", False),
    ("Professional", "$249", "/month", "Up to 25 users", "15 active projects",
     "+ Financial + Equipment + CRM", "CSV + Excel import", "Priority support", True),
    ("Enterprise", "$499", "/month", "Unlimited users", "Unlimited projects",
     "+ Property + AI + Portals", "+ P6/MSP + QuickBooks", "Dedicated CSM", False),
]

for i, (name, price, period, users, projects, modules, imports, support, featured) in enumerate(tiers):
    left = Inches(0.8 + i * 4.1)
    top = Inches(1.9)

    # Card background
    border_color = BLUE if featured else RGBColor(0xD6, 0xD3, 0xD1)
    card_bg = WHITE
    card = add_shape(slide, left, top, Inches(3.6), Inches(5.0), card_bg)
    card.line.color.rgb = border_color
    card.line.width = Pt(2 if featured else 1)

    if featured:
        # "Most Popular" badge
        badge = add_shape(slide, left + Inches(0.8), top - Inches(0.15), Inches(2.0), Inches(0.35), BLUE)
        add_text_box(slide, left + Inches(0.8), top - Inches(0.15), Inches(2.0), Inches(0.35),
                     "MOST POPULAR", font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Tier name
    add_text_box(slide, left + Inches(0.2), top + Inches(0.3), Inches(3.2), Inches(0.4),
                 name, font_size=20, color=BLUE if featured else TEXT_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=HEADING_FONT)

    # Price
    add_text_box(slide, left + Inches(0.2), top + Inches(0.8), Inches(3.2), Inches(0.6),
                 price, font_size=40, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.35), Inches(3.2), Inches(0.3),
                 period, font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Amber divider
    add_shape(slide, left + Inches(1.3), top + Inches(1.7), Inches(1.0), Pt(2), AMBER)

    # Features
    features = [users, projects, modules, imports, support]
    for j, feat in enumerate(features):
        add_text_box(slide, left + Inches(0.3), top + Inches(2.0 + j * 0.45), Inches(3.0), Inches(0.35),
                     f"\u2713  {feat}", font_size=12, color=TEXT_DARK)

# Unit economics below
add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
             "ARPA: $250/mo  |  LTV: $9,000  |  CAC: $500-$1,500  |  LTV:CAC: 6x-18x  |  Gross Margin: 85%+",
             font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 7)

# ═══════════════════════════════════════════════════════════
# SLIDE 8: Traction & What's Built
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_header_bar(slide)
add_logo(slide)

add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
             "What's built today",
             font_size=32, font_name=HEADING_FONT, color=TEXT_DARK, bold=True)

add_amber_bar(slide, Inches(0.6), Inches(1.5), Inches(1.0))

# Stats grid - 3 columns x 4 rows
stats_data = [
    ("80,000+", "Lines of Code", "TypeScript + SQL + CSS"),
    ("646", "Source Files", "100% TypeScript"),
    ("177", "API Endpoints", "REST architecture"),
    ("178", "Application Pages", "Full app router"),
    ("46", "Database Tables", "Full relational schema"),
    ("461", "Security Policies", "Row-Level Security"),
    ("35", "DB Migrations", "Progressive evolution"),
    ("9", "AI Providers", "30+ models supported"),
    ("13", "Demo Datasets", "Stadium to residential"),
    ("60", "Unit Tests", "Financial engine covered"),
    ("19", "Feature Modules", "PM to AI to property"),
    ("3", "User Portals", "Tenant, vendor, employee"),
]

for i, (number, label, detail) in enumerate(stats_data):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.2)
    top = Inches(1.9 + row * 1.25)

    card = add_shape(slide, left, top, Inches(3.8), Inches(1.05), BG_WARM)
    card.line.fill.background()

    add_text_box(slide, left + Inches(0.2), top + Inches(0.05), Inches(1.5), Inches(0.6),
                 number, font_size=32, color=BLUE, bold=True, font_name=HEADING_FONT)
    add_text_box(slide, left + Inches(1.8), top + Inches(0.1), Inches(1.8), Inches(0.3),
                 label, font_size=14, color=TEXT_DARK, bold=True)
    add_text_box(slide, left + Inches(1.8), top + Inches(0.4), Inches(1.8), Inches(0.3),
                 detail, font_size=11, color=MUTED)

add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
             "Development equivalent: 12-18 months of a 3-4 person engineering team",
             font_size=14, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 8)

# ═══════════════════════════════════════════════════════════
# SLIDE 9: Financial Projections
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_header_bar(slide)
add_logo(slide)

add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
             "Path to $1M ARR in 24 months",
             font_size=32, font_name=HEADING_FONT, color=TEXT_DARK, bold=True)

add_amber_bar(slide, Inches(0.6), Inches(1.5), Inches(1.0))

# Revenue timeline bars
quarters = [
    ("Q2 '26", 5, "$1.2K", "$15K"),
    ("Q3 '26", 15, "$3.8K", "$45K"),
    ("Q4 '26", 35, "$8.3K", "$99K"),
    ("Q1 '27", 60, "$15K", "$180K"),
    ("Q2 '27", 100, "$25K", "$300K"),
    ("Q3 '27", 150, "$38K", "$450K"),
    ("Q4 '27", 250, "$63K", "$750K"),
    ("Q1 '28", 350, "$88K", "$1.05M"),
]

max_mrr = 88000
for i, (qtr, customers, mrr, arr) in enumerate(quarters):
    top = Inches(1.9 + i * 0.6)
    # Quarter label
    add_text_box(slide, Inches(0.6), top, Inches(0.9), Inches(0.35),
                 qtr, font_size=12, color=TEXT_DARK, bold=True)

    # Bar
    bar_max = Inches(7.0)
    mrr_val = int(mrr.replace("$", "").replace("K", "000").replace(".", "").replace("M", "00000")
                   .replace(",", ""))
    # Approximate the numeric value
    bar_width = max(Inches(0.3), Inches(7.0 * min(1.0, (i + 1) / 8)))
    color = BLUE if i < 7 else GREEN
    add_shape(slide, Inches(1.7), top + Inches(0.05), bar_width, Inches(0.3), color)

    # MRR label on bar
    add_text_box(slide, Inches(1.8), top + Inches(0.02), Inches(2), Inches(0.3),
                 f"MRR: {mrr}", font_size=11, color=WHITE, bold=True)

    # Customers & ARR
    add_text_box(slide, Inches(9.5), top, Inches(1.5), Inches(0.3),
                 f"{customers} customers", font_size=11, color=TEXT_DARK)
    add_text_box(slide, Inches(11.2), top, Inches(1.5), Inches(0.3),
                 f"ARR: {arr}", font_size=11, color=BLUE, bold=True)

# Milestones
add_shape(slide, Inches(0.6), Inches(6.8), Inches(12.0), Inches(0.5), BG_WARM)
add_text_box(slide, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.4),
             "Assumptions:  ARPA $250/mo  |  5% monthly growth  |  <3% churn  |  CAC payback 3-6 months",
             font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 9)

# ═══════════════════════════════════════════════════════════
# SLIDE 10: The Ask
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_header_bar(slide)
add_logo(slide)

add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
             "Pre-Seed: $750K at $4M-$6M pre-money",
             font_size=32, font_name=HEADING_FONT, color=TEXT_DARK, bold=True)

add_amber_bar(slide, Inches(0.6), Inches(1.5), Inches(1.0))

# Use of Funds pie-chart substitute (horizontal bars)
funds = [
    ("Sales & Marketing", "$300K", "40%", "First 50 customers, content, trade shows"),
    ("Engineering", "$225K", "30%", "Test coverage, QuickBooks sync, features"),
    ("Operations", "$150K", "20%", "Infrastructure, support, legal"),
    ("Reserve", "$75K", "10%", "Working capital buffer"),
]

colors = [BLUE, RGBColor(0x3B, 0x82, 0xF6), AMBER, MUTED]

for i, (category, amount, pct, desc) in enumerate(funds):
    top = Inches(1.9 + i * 1.1)

    # Color bar
    bar_w = Inches(float(pct.replace("%", "")) / 100 * 8)
    add_shape(slide, Inches(0.6), top, bar_w, Inches(0.5), colors[i])

    add_text_box(slide, Inches(0.8), top + Inches(0.05), Inches(4), Inches(0.4),
                 f"{category}  —  {amount} ({pct})", font_size=16, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), top + Inches(0.55), Inches(8), Inches(0.3),
                 desc, font_size=13, color=MUTED)

# Milestones this round enables
add_text_box(slide, Inches(0.6), Inches(5.7), Inches(12), Inches(0.4),
             "Milestones this round enables:", font_size=18, color=TEXT_DARK, bold=True)

milestones = [
    "First 50 paying customers ($12K+ MRR)",
    "Product-market fit validation",
    "Complete QuickBooks integration",
    "Seed-ready metrics for $2M-$5M raise",
]

for i, m in enumerate(milestones):
    add_text_box(slide, Inches(0.8), Inches(6.2 + i * 0.35), Inches(11), Inches(0.3),
                 f"\u2713  {m}", font_size=13, color=TEXT_DARK)

add_slide_number(slide, 10)

# ═══════════════════════════════════════════════════════════
# SLIDE 11: Why Now
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_header_bar(slide)
add_logo(slide)

add_text_box(slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
             "Four converging tailwinds",
             font_size=32, font_name=HEADING_FONT, color=TEXT_DARK, bold=True)

add_amber_bar(slide, Inches(0.6), Inches(1.5), Inches(1.0))

tailwinds = [
    ("AI Maturity", "LLMs can now understand construction-specific context. Buildwrk is the first platform to put AI on real construction data with function calling.", "\U0001F916"),
    ("Cloud Adoption", "COVID forced construction to accept cloud software. The laggard industry is now the fastest-growing SaaS adoption segment.", "\u2601\ufe0f"),
    ("Infrastructure Spending", "The $1.2T Infrastructure Investment and Jobs Act creates thousands of new projects needing management software.", "\U0001F3D7\ufe0f"),
    ("Generational Shift", "Millennial PMs and superintendents expect modern, mobile-first software. They won't tolerate 1990s interfaces.", "\U0001F4F1"),
]

for i, (title, desc, icon) in enumerate(tailwinds):
    col = i % 2
    row = i // 2
    left = Inches(0.6 + col * 6.3)
    top = Inches(2.0 + row * 2.5)

    card = add_shape(slide, left, top, Inches(5.8), Inches(2.1), BG_WARM)
    card.line.fill.background()

    # Number
    add_shape(slide, left + Inches(0.2), top + Inches(0.2), Inches(0.5), Inches(0.5), BLUE)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(0.5), Inches(0.5),
                 str(i + 1), font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(0.9), top + Inches(0.2), Inches(4.6), Inches(0.4),
                 title, font_size=18, color=TEXT_DARK, bold=True)
    add_text_box(slide, left + Inches(0.9), top + Inches(0.7), Inches(4.6), Inches(1.2),
                 desc, font_size=13, color=MUTED)

add_slide_number(slide, 11)

# ═══════════════════════════════════════════════════════════
# SLIDE 12: Contact / Closing
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Full blue left panel
add_shape(slide, Inches(0), Inches(0), Inches(5.5), SLIDE_H, BLUE)

# Logo on blue
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(4.2), Inches(1.0),
             "Buildwrk", font_size=48, font_name=HEADING_FONT, color=WHITE, bold=True)

add_amber_bar(slide, Inches(0.8), Inches(2.6), Inches(1.2), Pt(4))

add_multi_text(slide, Inches(0.8), Inches(2.8), Inches(4.2), Inches(3.0), [
    ("The Unified Construction", 24, WHITE, False, HEADING_FONT),
    ("Intelligence Platform", 24, WHITE, False, HEADING_FONT),
    ("", 14, WHITE),
    ('"Every hour spent reconciling data', 13, RGBColor(0xD1, 0xD5, 0xDB), False),
    ("between disconnected tools is an hour", 13, RGBColor(0xD1, 0xD5, 0xDB), False),
    ('not spent building."', 13, RGBColor(0xD1, 0xD5, 0xDB), False),
])

# Right side - contact info
add_text_box(slide, Inches(6.2), Inches(2.0), Inches(6.5), Inches(0.5),
             "Let's Talk", font_size=36, color=TEXT_DARK, bold=True, font_name=HEADING_FONT)

add_amber_bar(slide, Inches(6.2), Inches(2.65), Inches(0.8))

contact_items = [
    ("Live Demo", "construction-gamma-six.vercel.app"),
    ("Email", "[your-email@domain.com]"),
    ("Phone", "[your-phone-number]"),
    ("Location", "[your-city, state]"),
]

for i, (label, value) in enumerate(contact_items):
    top = Inches(3.0 + i * 0.7)
    add_text_box(slide, Inches(6.2), top, Inches(2), Inches(0.3),
                 label, font_size=13, color=MUTED, bold=True)
    add_text_box(slide, Inches(8.2), top, Inches(4.5), Inches(0.3),
                 value, font_size=15, color=TEXT_DARK)

add_text_box(slide, Inches(6.2), Inches(5.8), Inches(6.5), Inches(0.4),
             "This pitch deck is confidential.", font_size=11, color=MUTED)

add_slide_number(slide, 12)


# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
output_dir = r"c:\Users\beltr\Construction.erp\docs\business"
output_path = os.path.join(output_dir, "Buildwrk_Pitch_Deck.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
